import json
import sys
from pathlib import Path
import io
from PIL import Image
import numpy as np
import pymupdf  # PyMuPDF
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from ocr_engine import OCREngineFactory, OCREngine

class FileConverter:
    def __init__(self, config_path: str = "config.json"):
        # Resolve config_path relative to this script's directory
        base_dir = Path(__file__).parent.absolute()
        
        # Try multiple locations for config.json (consistent with RagEngine)
        config_locations = [
            base_dir / config_path,
            Path(config_path),
            base_dir / "config.json"
        ]
        
        self.config = None
        for loc in config_locations:
            if loc.exists():
                try:
                    with open(loc, "r", encoding="utf-8") as f:
                        self.config = json.load(f)
                    break
                except Exception as e:
                    print(f"Error loading config from {loc}: {e}", file=sys.stderr)

        if self.config is None:
            # Fallback to empty dict if config not found, but warn
            print(f"Warning: Could not find config file. Using defaults.", file=sys.stderr)
            self.config = {}

        # Additional text extensions from config
        self.extra_text_extensions = self.config.get("extra_text_extensions", [])
        
        # OCRエンジンの初期化（抽象化レイヤーを使用）
        self.ocr_engine_instance: OCREngine = OCREngineFactory.create_engine(self.config)
        print(f"OCR engine initialized: {self.ocr_engine_instance.get_name()}", file=sys.stderr)

    def perform_ocr_from_bytes(self, image_data: bytes) -> str:
        """メモリ上の画像データからテキストを抽出する"""
        return self.ocr_engine_instance.perform_ocr_from_bytes(image_data)

    def perform_ocr(self, image_path: Path) -> str:
        """ファイルパスからテキストを抽出する"""
        return self.ocr_engine_instance.perform_ocr(image_path)

    def extract_text_from_pdf(self, pdf_path: Path, progress_callback=None) -> str:
        try:
            doc = pymupdf.open(pdf_path)
            pages_text = []
            total_pages = len(doc)
            for i, page in enumerate(doc):
                pix = page.get_pixmap(dpi=300)
                img_data = pix.tobytes("png")
                print(f"Converting PDF page {i+1}/{total_pages}...", file=sys.stderr)
                if progress_callback:
                    progress_callback(i + 1, total_pages, pdf_path.name)
                page_text = self.perform_ocr_from_bytes(img_data)
                pages_text.append((i + 1, page_text))
            doc.close()

            # 全ページでOCRテキストが取得できなかった場合はファイル生成しない
            if not any(text.strip() for _, text in pages_text):
                print(f"OCR failed for all pages in {pdf_path.name}. Skipping file generation.", file=sys.stderr)
                return ""

            full_text = f"# PDF Content: {pdf_path.name}\n\n"
            for page_num, page_text in pages_text:
                full_text += f"\n## Page {page_num}\n\n{page_text}\n"
            return full_text
        except Exception as e:
            print(f"Error reading PDF {pdf_path}: {e}", file=sys.stderr)
            return ""

    def extract_text_from_docx(self, file_path: Path) -> str:
        try:
            doc = Document(file_path)
            full_text = [f"# DOCX Content: {file_path.name}\n"]
            for para in doc.paragraphs:
                full_text.append(para.text)
            return "\n\n".join(full_text)
        except Exception as e:
            print(f"Error reading DOCX {file_path}: {e}", file=sys.stderr)
            return ""

    def extract_text_from_xlsx(self, file_path: Path) -> str:
        # BUG-014 fix: Try data_only=True first, fallback to data_only=False for formulas
        for data_only in [True, False]:
            try:
                wb = load_workbook(file_path, data_only=data_only)
                text_parts = [f"# XLSX Content: {file_path.name}\n"]
                for sheet in wb.sheetnames:
                    ws = wb[sheet]
                    text_parts.append(f"## Sheet: {sheet}")
                    # Simple Markdown table conversion could be added here, currently tab-separated
                    for row in ws.iter_rows(values_only=True):
                        row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                        if row_text.strip().replace("|", ""):
                             text_parts.append(f"| {row_text} |")
                return "\n\n".join(text_parts)
            except Exception as e:
                if data_only:
                    print(f"Warning: data_only=True failed for {file_path.name}, trying data_only=False: {e}", file=sys.stderr)
                else:
                    print(f"Error reading XLSX {file_path}: {e}", file=sys.stderr)
        return ""

    def extract_text_from_pptx(self, file_path: Path) -> str:
        try:
            prs = Presentation(file_path)
            text_parts = [f"# PPTX Content: {file_path.name}\n"]
            for i, slide in enumerate(prs.slides):
                text_parts.append(f"## Slide {i+1}")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)
            return "\n\n".join(text_parts)
        except Exception as e:
            print(f"Error reading PPTX {file_path}: {e}", file=sys.stderr)
            return ""

    def convert_file(self, file_path: Path, output_path: Path, progress_callback=None):
        """ファイルを変換してMarkdownとして保存する。
        戻り値:
            True  - 変換・保存に成功
            False - スキップ（未サポート形式・タイムスタンプが新しいので変換不要）
            None  - 変換を試みたがOCR失敗等でコンテンツが取得できなかった
        """
        if file_path.name.startswith("~$"):
            return False
            
        if output_path.exists():
            if output_path.stat().st_mtime > file_path.stat().st_mtime:
                print(f"Skipping {file_path.name} (uptodate)", file=sys.stderr)
                return False

        print(f"Converting {file_path.name}...", file=sys.stderr)
        attempted = False  # 変換を試みたかどうか
        content = None
        if file_path.suffix.lower() == ".pdf":
            attempted = True
            content = self.extract_text_from_pdf(file_path, progress_callback=progress_callback)
        elif file_path.suffix.lower() == ".docx":
            attempted = True
            content = self.extract_text_from_docx(file_path)
        elif file_path.suffix.lower() == ".xlsx":
            attempted = True
            content = self.extract_text_from_xlsx(file_path)
        elif file_path.suffix.lower() == ".pptx":
            attempted = True
            content = self.extract_text_from_pptx(file_path)
        elif file_path.suffix.lower() in [".png", ".jpg", ".jpeg", ".bmp"]:
            attempted = True
            content = self.perform_ocr(file_path)
        elif file_path.suffix.lower() in [
            # テキスト・マークダウン
            ".txt", ".md",
            # Web
            ".html", ".htm", ".css", ".js", ".ts",
            # データ形式
            ".json", ".xml", ".yaml", ".yml", ".toml", ".ini", ".env",
            # Python
            ".py",
            # C / C++
            ".c", ".cpp", ".cc", ".cxx", ".h", ".hpp", ".hxx",
            # Verilog / SystemVerilog
            ".v", ".sv", ".svh", ".vh",
            # Rust
            ".rs",
            # Go
            ".go",
            # Java / Kotlin
            ".java", ".kt",
            # その他
            ".sh", ".bash", ".bat", ".ps1", ".rb", ".php", ".swift", ".cs",
        ] + self.extra_text_extensions:
            try:
                with open(file_path, "r", encoding="utf-8") as f:
                    content = f"# Source: {file_path.name}\n\n" + f.read()
            except Exception:
                content = None
        
        if content:
            output_path.parent.mkdir(parents=True, exist_ok=True)
            with open(output_path, "w", encoding="utf-8") as f:
                f.write(content)
            print(f"Saved to {output_path}", file=sys.stderr)
            return True
        elif attempted:
            # 変換を試みたが内容が空（OCR失敗等）
            print(f"OCR/conversion produced no content for {file_path.name} (skipping file generation)", file=sys.stderr)
            return None
        else:
            print(f"Skipped {file_path.name} (unsupported or error)", file=sys.stderr)
            return False

def main():
    converter = FileConverter()
    
    # Input/Output directories
    base_dir = Path(__file__).parent.absolute()
    
    source_dir_name = converter.config.get("source_docs_dir", "documents")
    docs_dir = (base_dir / source_dir_name).resolve()
    
    output_dir_name = converter.config.get("docs_dir", "converted_docs")
    output_dir = (base_dir / output_dir_name).resolve()
    output_dir.mkdir(parents=True, exist_ok=True)
    
    print(f"Scanning {docs_dir}...", file=sys.stderr)
    if not docs_dir.exists():
        print(f"Directory {docs_dir} not found.", file=sys.stderr)
        return

    # 1. Sync: Convert new/modified files
    valid_markdown_files = set()
    
    for file_path in docs_dir.rglob("*"):
        if file_path.is_file():
            # Calculate output path preserving subdirectory structure
            rel_path = file_path.relative_to(docs_dir)
            # Replace extension with .md for the output file
            # If original file is .md, avoiding conflict is good, but typically we map .pdf -> .md
            # If we have foo.pdf, we get foo.md.
            output_rel_path = rel_path.with_suffix(".md")
            output_path = output_dir / output_rel_path
            
            # Convert (will check timestamp inside)
            converter.convert_file(file_path, output_path)
            
            # Keep track of valid output files to detect orphans
            valid_markdown_files.add(output_path.resolve())

    # 2. Cleanup: Remove orphaned markdown files in converted_docs
    # Only iterate if we effectively scanned documents.
    print(f"Checking for deleted files in {output_dir}...", file=sys.stderr)
    for md_file in output_dir.rglob("*.md"):
        if md_file.is_file():
            if md_file.resolve() not in valid_markdown_files:
                print(f"Removing orphaned file: {md_file}", file=sys.stderr)
                try:
                    md_file.unlink()
                    # Try to remove empty parent directories
                    parent = md_file.parent
                    while parent != output_dir:
                        if not any(parent.iterdir()):
                            parent.rmdir()
                            parent = parent.parent
                        else:
                            break
                except Exception as e:
                    print(f"Error removing {md_file}: {e}", file=sys.stderr)

    print("File conversion completed.", file=sys.stderr)

if __name__ == "__main__":
    main()
